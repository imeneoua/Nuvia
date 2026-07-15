from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Create presentation
prs = Presentation()

# Slide 1: Title
slide_layout = prs.slide_layouts[0] # Title slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "NUVIA"
p = slide.placeholders[1]
p.text = "L'Écosystème Culinaire Nouvelle Génération\nDe l'inspiration à l'assiette."

# Helper function to add a content slide
def add_slide(title_text, bullet_points):
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    title = slide.shapes.title
    title.text = title_text
    
    body = slide.placeholders[1]
    tf = body.text_frame
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point
            
    # Set font size for all text
    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(28)

# Slide 2
add_slide(
    "Le Problème",
    [
        "Le manque d'inspiration au quotidien.",
        "La désorganisation (gaspillage, courses mal gérées).",
        "La barrière de la langue sur les recettes mondiales.",
        "Le manque d'assistance en temps réel pour cuisiner."
    ]
)

# Slide 3
add_slide(
    "La Solution : La Plateforme Nuvia",
    [
        "Explore : Une infinité de recettes à découvrir.",
        "Planner : Planification hebdomadaire intelligente.",
        "Auth & Profile : Un espace sécurisé, connecté via Google OAuth.",
        "Dashboard Admin : Panneau de contrôle complet pour gérer les utilisateurs."
    ]
)

# Slide 4
add_slide(
    "La Magie de l'IA (L'Assistant Intégré)",
    [
        "Plus qu'un simple site web : Un véritable compagnon.",
        "Mode Chat & Vocal immersif (Mains-libres).",
        "RAG (Retrieval-Augmented Generation) sur +1 Million de recettes locales.",
        "Traduction instantanée à la volée (FR, EN, AR) via l'IA."
    ]
)

# Slide 5
add_slide(
    "L'Envers du Décor (Architecture Technique)",
    [
        "Front-end : React + Vite (Design premium Glassmorphism).",
        "Back-end : FastAPI (Python) ultra-rapide.",
        "Base de données : SQLite pour la gestion sécurisée des utilisateurs.",
        "LLM & Voix : Gemini 1.5 Flash + Web Speech API (Reconnaissance vocale)."
    ]
)

# Slide 6
add_slide(
    "Démonstration",
    [
        "1. Fluidité de la connexion Google.",
        "2. Découverte de la plateforme (Planner et Explore).",
        "3. Posez une question au Chatbot vocal (en français ou arabe).",
        "4. Accès au portail Administrateur."
    ]
)

# Slide 7
add_slide(
    "Conclusion & Perspectives",
    [
        "Un écosystème moderne prêt pour la production.",
        "Évolution prévue vers une App Mobile native.",
        "Génération automatique de listes de courses avec des partenaires locaux.",
        "",
        "Merci pour votre attention !"
    ]
)

# Save to the workspace
output_path = os.path.join(os.path.dirname(__file__), "Nuvia_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to {output_path}")
